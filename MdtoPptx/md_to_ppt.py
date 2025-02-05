import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt

def parse_markdown(file_path):
    slides = []
    current_slide = None

    # Regular expression to match markdown bullet items
    # Captures the indent (spaces) and the text after the dash.
    bullet_pattern = re.compile(r'^( *)(-)\s+(.*)$')

    with open(file_path, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.rstrip()
            if not line:
                continue  # skip empty lines
            match = bullet_pattern.match(line)
            if match:
                indent, dash, content = match.groups()
                indent_level = len(indent)
                # Consider indent_level == 0 as a new slide
                if indent_level == 0:
                    # Start a new slide with this bullet as title
                    current_slide = {'title': content, 'bullets': []}
                    slides.append(current_slide)
                else:
                    # Nested bullet: add to the current slide if one exists
                    if current_slide is not None:
                        current_slide['bullets'].append(content)
            else:
                # If the line doesn't match the bullet pattern, you might decide to add it
                # as plain text. Here, we ignore non-bullet lines.
                continue

    return slides

def create_pptx(slides, output_path):
    prs = Presentation()
    # Use a title and content layout (usually layout index 1)
    slide_layout = prs.slide_layouts[1]

    for slide_data in slides:
        slide = prs.slides.add_slide(slide_layout)
        # Set the title text
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_data['title']

        # Set the bullet points in the content placeholder if there are any
        if slide_data['bullets']:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear any default text
            for bullet in slide_data['bullets']:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 1  # nested bullet
                p.font.size = Pt(18)
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

def main():
    if len(sys.argv) != 3:
        print("Usage: python md_to_ppt.py <input_markdown_file> <output_pptx_file>")
        sys.exit(1)
    input_md = sys.argv[1]
    output_pptx = sys.argv[2]
    slides = parse_markdown(input_md)
    create_pptx(slides, output_pptx)

if __name__ == '__main__':
    main()
